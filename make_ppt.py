from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0b, 0x11, 0x20)   # #0b1120
PURPLE    = RGBColor(0x66, 0x7e, 0xea)   # #667eea
MAGENTA   = RGBColor(0xf0, 0x93, 0xfb)   # #f093fb
WHITE     = RGBColor(0xff, 0xff, 0xff)
MUTED     = RGBColor(0x94, 0xa3, 0xb8)   # slate-400
GREEN     = RGBColor(0x10, 0xb9, 0x81)   # verified
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)   # unverified
RED       = RGBColor(0xef, 0x44, 0x44)   # hallucinated
CARD_BG   = RGBColor(0x1e, 0x2a, 0x40)   # card surface

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def bg(slide):
    add_rect(slide, 0, 0, W, H, fill=BG_DARK)

def gradient_bar(slide, y=Inches(0.08), h=Inches(0.06)):
    """Top accent bar purple→magenta"""
    add_rect(slide, 0, y, W/2, h, fill=PURPLE)
    add_rect(slide, W/2, y, W/2, h, fill=MAGENTA)

def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}",
             W - Inches(1.2), H - Inches(0.4), Inches(1), Inches(0.3),
             size=Pt(10), color=MUTED, align=PP_ALIGN.RIGHT)

def section_tag(slide, label, color=PURPLE):
    tag = add_rect(slide, Inches(0.5), Inches(0.28), Inches(2.4), Inches(0.32),
                   fill=RGBColor(int(color[0]*0.15), int(color[1]*0.15), int(color[2]*0.15)))
    add_text(slide, label, Inches(0.52), Inches(0.28), Inches(2.4), Inches(0.32),
             size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body, icon="", title_color=WHITE):
    add_rect(slide, x, y, w, h, fill=CARD_BG,
             line=RGBColor(0x2d, 0x3d, 0x5a), line_w=Pt(1))
    add_text(slide, f"{icon}  {title}" if icon else title,
             x+Inches(0.18), y+Inches(0.15), w-Inches(0.3), Inches(0.38),
             size=Pt(13), bold=True, color=title_color)
    add_text(slide, body,
             x+Inches(0.18), y+Inches(0.55), w-Inches(0.3), h-Inches(0.65),
             size=Pt(11), color=MUTED)

TOTAL = 10

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
# big logo text
add_text(s, "TruthTrace", Inches(1), Inches(1.8), Inches(11), Inches(1.6),
         size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "The AI Hallucination Auditor",
         Inches(1), Inches(3.3), Inches(11), Inches(0.7),
         size=Pt(26), bold=False, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(s, "Paste any AI-generated text  →  Every fact scraped  →  You decide the truth",
         Inches(1), Inches(4.1), Inches(11), Inches(0.6),
         size=Pt(16), italic=True, color=MUTED, align=PP_ALIGN.CENTER)
# badges
for i, label in enumerate(["Real-time Scraping", "LLM-as-Judge", "User-Controlled Verdict", "0–100 Trust Score"]):
    bx = Inches(1.2 + i * 2.7)
    add_rect(s, bx, Inches(5.2), Inches(2.4), Inches(0.42),
             fill=RGBColor(0x1e, 0x2a, 0x40), line=PURPLE, line_w=Pt(1))
    add_text(s, label, bx, Inches(5.2), Inches(2.4), Inches(0.42),
             size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "truthtrace.vercel.app", Inches(1), Inches(6.6), Inches(11), Inches(0.4),
         size=Pt(13), color=MAGENTA, align=PP_ALIGN.CENTER)
slide_number(s, 1, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "THE PROBLEM", RED)
add_text(s, "AI Hallucinations Are Everywhere",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.8),
         size=Pt(36), bold=True, color=WHITE)
add_text(s, "AI models confidently generate false facts — names, dates, papers, statistics — that never existed.",
         Inches(0.5), Inches(1.55), Inches(12), Inches(0.6),
         size=Pt(16), color=MUTED)

examples = [
    ("❌  Invented Person", '"Dr. Harold Venkatesh of MIT published a landmark study in Nature Structural Engineering" — this person does not exist.'),
    ("❌  Wrong Number", '"The Eiffel Tower is 450 metres tall" — it is 330 metres.'),
    ("❌  Fake Citation", '"Einstein won the Nobel Prize for Relativity" — he won it for the photoelectric effect.'),
    ("❌  Made-up Event", '"In 2019 the tower was painted by Saint-Gobain" — no such record exists.'),
]
for i, (title, body) in enumerate(examples):
    row, col = divmod(i, 2)
    card(s, Inches(0.4 + col*6.4), Inches(2.35 + row*2.2), Inches(6.1), Inches(2.0),
         title, body, title_color=RED)
slide_number(s, 2, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT IS TRUTHTRACE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "SOLUTION", GREEN)
add_text(s, "What is TruthTrace?",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

add_text(s,
    "TruthTrace is an AI-powered hallucination auditor. You paste any AI-generated document, "
    "and TruthTrace automatically:\n\n"
    "  1.  Extracts every verifiable factual claim\n"
    "  2.  Scrapes Wikipedia + DuckDuckGo for real-world evidence\n"
    "  3.  Shows you a source agreement score (% of sources that agree)\n"
    "  4.  Lets YOU decide the final verdict — Verified, Unverified, or Hallucinated\n"
    "  5.  Computes a 0–100 Trust Score based on your decisions",
    Inches(0.6), Inches(1.6), Inches(7.5), Inches(5.2),
    size=Pt(16), color=WHITE)

# Right side quote box
add_rect(s, Inches(8.5), Inches(1.6), Inches(4.3), Inches(5.2),
         fill=CARD_BG, line=PURPLE, line_w=Pt(1.5))
add_text(s, "Think of it as...",
         Inches(8.7), Inches(1.8), Inches(3.9), Inches(0.4),
         size=Pt(13), bold=True, color=PURPLE)
add_text(s,
    '"Spell-check — but\nfor facts, not spelling."',
    Inches(8.7), Inches(2.3), Inches(3.9), Inches(1.4),
    size=Pt(20), bold=True, italic=True, color=WHITE)
add_text(s,
    "Every factual claim gets individually verified against live web sources. "
    "No more blindly trusting AI-generated content.",
    Inches(8.7), Inches(4.0), Inches(3.9), Inches(2.5),
    size=Pt(13), color=MUTED)
slide_number(s, 3, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS (pipeline)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "HOW IT WORKS", PURPLE)
add_text(s, "The 5-Step Pipeline",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

steps = [
    ("1", "Paste Document", "User pastes any AI-generated text into the app"),
    ("2", "Extract Claims", "Groq LLM reads the text and pulls out every atomic, verifiable fact"),
    ("3", "Scrape Evidence", "Wikipedia REST API + DuckDuckGo API scraped for each claim — up to 15 sources"),
    ("4", "Score Evidence", "AI counts how many sources agree vs disagree → gives a % score"),
    ("5", "User Decides", "You see the evidence and click Verified / Unverified / Hallucinated"),
]

for i, (num, title, body) in enumerate(steps):
    x = Inches(0.3 + i * 2.55)
    # Number circle
    add_rect(s, x + Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.8),
             fill=PURPLE, line=None)
    add_text(s, num, x + Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.8),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow (except last)
    if i < 4:
        add_text(s, "→", x + Inches(1.6), Inches(1.75), Inches(0.9), Inches(0.6),
                 size=Pt(20), color=PURPLE, align=PP_ALIGN.CENTER)
    # Card
    add_rect(s, x, Inches(2.65), Inches(2.35), Inches(3.8),
             fill=CARD_BG, line=RGBColor(0x2d, 0x3d, 0x5a), line_w=Pt(1))
    add_text(s, title, x+Inches(0.12), Inches(2.8), Inches(2.1), Inches(0.45),
             size=Pt(13), bold=True, color=PURPLE)
    add_text(s, body, x+Inches(0.12), Inches(3.3), Inches(2.1), Inches(2.8),
             size=Pt(12), color=MUTED)
slide_number(s, 4, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EVIDENCE SCRAPING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "EVIDENCE ENGINE", PURPLE)
add_text(s, "How Evidence is Scraped",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

sources = [
    ("📖  Wikipedia REST API", "Fetches the direct article summary for the claim topic. Also searches for 3 related Wikipedia articles to cross-reference.", PURPLE),
    ("🔍  DuckDuckGo Instant API", "Hits DuckDuckGo's free JSON API — returns abstract text, related topics, and infobox facts (dates, statistics, names).", MAGENTA),
    ("🌐  Real Web Page Fetching", "Visits top 5 URLs from search results, strips HTML, and extracts the main body text to read actual source content.", GREEN),
    ("🤖  AI Evidence Analyst", "Groq AI reads all sources and counts: how many support the claim? How many contradict? Returns a neutral summary — NO verdict.", AMBER),
]
for i, (title, body, color) in enumerate(sources):
    row, col = divmod(i, 2)
    cx = Inches(0.4 + col * 6.4)
    cy = Inches(1.7 + row * 2.5)
    add_rect(s, cx, cy, Inches(6.1), Inches(2.3),
             fill=CARD_BG, line=color, line_w=Pt(1.5))
    add_text(s, title, cx+Inches(0.18), cy+Inches(0.15), Inches(5.7), Inches(0.4),
             size=Pt(13), bold=True, color=color)
    add_text(s, body, cx+Inches(0.18), cy+Inches(0.6), Inches(5.7), Inches(1.5),
             size=Pt(12), color=MUTED)
add_text(s, "Result: up to 15 sources per claim, all scraped live in real-time",
         Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
slide_number(s, 5, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — USER CONTROLS VERDICT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "KEY DESIGN DECISION", MAGENTA)
add_text(s, "You Decide — Not the AI",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)
add_text(s, "The AI only presents evidence. The final verdict is always yours.",
         Inches(0.5), Inches(1.45), Inches(12), Inches(0.45),
         size=Pt(16), color=MUTED)

# Old vs New
for col, (label, items, color) in enumerate([
    ("❌  Old Way (AI decides)", [
        "AI labels everything Verified / Hallucinated",
        "You have no say in the outcome",
        "If AI is wrong, you're misled",
        "Black box — no evidence shown",
    ], RED),
    ("✅  TruthTrace (You decide)", [
        "AI only shows: X of Y sources agree",
        "You see all sources and snippets",
        "You click Verified / Unverified / Hallucinated",
        "Trust score built from YOUR decisions only",
    ], GREEN),
]):
    cx = Inches(0.5 + col * 6.4)
    add_rect(s, cx, Inches(2.1), Inches(6.0), Inches(4.6),
             fill=CARD_BG, line=color, line_w=Pt(2))
    add_text(s, label, cx+Inches(0.2), Inches(2.25), Inches(5.6), Inches(0.5),
             size=Pt(14), bold=True, color=color)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", cx+Inches(0.2), Inches(2.9 + j*0.85), Inches(5.6), Inches(0.7),
                 size=Pt(13), color=WHITE if col else MUTED)
slide_number(s, 6, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRUST SCORE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "TRUST SCORE", GREEN)
add_text(s, "The 0–100 Trust Score",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

add_text(s, "Formula:", Inches(0.6), Inches(1.6), Inches(5), Inches(0.4),
         size=Pt(14), bold=True, color=PURPLE)

# formula box
add_rect(s, Inches(0.6), Inches(2.1), Inches(7), Inches(0.85),
         fill=CARD_BG, line=PURPLE, line_w=Pt(1.5))
add_text(s, "Trust Score  =  (Verified by you / Total claims)  ×  Avg source score of verified claims",
         Inches(0.8), Inches(2.15), Inches(6.6), Inches(0.75),
         size=Pt(13), bold=True, color=WHITE)

add_text(s, "What this means:",
         Inches(0.6), Inches(3.15), Inches(5), Inches(0.4),
         size=Pt(14), bold=True, color=PURPLE)
bullets = [
    "Score only updates when YOU label a claim — starts at 0",
    "Labelling more claims = more accurate score",
    "Even one hallucinated claim heavily penalises the score",
    "Score range:  0 = completely untrustworthy  →  100 = fully verified",
]
for i, b in enumerate(bullets):
    add_text(s, f"•  {b}", Inches(0.7), Inches(3.65 + i*0.7), Inches(7.5), Inches(0.6),
             size=Pt(14), color=WHITE)

# Score examples on right
for i, (score, label, color) in enumerate([
    ("87", "High Trust", GREEN),
    ("42", "Moderate Trust", AMBER),
    ("11", "Low Trust", RED),
]):
    bx = Inches(9.0)
    by = Inches(1.5 + i * 1.9)
    add_rect(s, bx, by, Inches(3.5), Inches(1.6),
             fill=CARD_BG, line=color, line_w=Pt(2))
    add_text(s, score, bx+Inches(0.2), by+Inches(0.1), Inches(1.2), Inches(1.3),
             size=Pt(52), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, label, bx+Inches(1.5), by+Inches(0.5), Inches(1.8), Inches(0.6),
             size=Pt(16), bold=True, color=WHITE)
slide_number(s, 7, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "TECH STACK", PURPLE)
add_text(s, "Built With",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

stack = [
    ("⚡  Framework",       "Next.js 14 App Router + TypeScript"),
    ("🎨  Styling",         "Tailwind CSS v4 + Custom CSS Design System"),
    ("✨  Animation",       "Framer Motion"),
    ("🤖  LLM",            "Groq SDK — Llama 3.3 70B / Llama 4 Scout"),
    ("🔍  Evidence",       "Wikipedia REST API + DuckDuckGo JSON API"),
    ("🕷️  Scraping",       "Node.js fetch — reads actual web pages"),
    ("📡  Streaming",      "Server-Sent Events (SSE) — live result feed"),
    ("📊  Charts",         "Recharts — Trust Score ring + Verdict bar"),
    ("☁️  Deploy",         "Vercel (free tier, zero config)"),
    ("🔒  Privacy",        "API key never stored — sent only at request time"),
]
for i, (layer, tech) in enumerate(stack):
    row, col = divmod(i, 2)
    cx = Inches(0.4 + col * 6.4)
    cy = Inches(1.7 + row * 1.08)
    add_rect(s, cx, cy, Inches(6.1), Inches(0.9),
             fill=CARD_BG, line=RGBColor(0x2d, 0x3d, 0x5a), line_w=Pt(1))
    add_text(s, layer, cx+Inches(0.15), cy+Inches(0.12), Inches(2.0), Inches(0.65),
             size=Pt(12), bold=True, color=PURPLE)
    add_text(s, tech, cx+Inches(2.0), cy+Inches(0.12), Inches(3.9), Inches(0.65),
             size=Pt(12), color=WHITE)
slide_number(s, 8, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
section_tag(s, "FEATURES", GREEN)
add_text(s, "What Makes TruthTrace Different",
         Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
         size=Pt(36), bold=True, color=WHITE)

features = [
    ("🔴  Live Streaming", "Results appear one by one as each claim is verified — no waiting for batch completion", GREEN),
    ("🕷️  Real Scraping", "Doesn't use cached data — scrapes Wikipedia + DuckDuckGo live for every audit", PURPLE),
    ("👤  User Verdict", "You control the outcome — AI shows evidence, you label each claim yourself", MAGENTA),
    ("📄  Annotated Doc", "Original text highlighted green/amber/red — hover any phrase to see the evidence", AMBER),
    ("📊  Trust Score Ring", "Live 0–100 gauge that updates every time you label a claim", GREEN),
    ("💾  Export Reports", "Download full audit as Markdown or JSON — includes all sources and scores", PURPLE),
]
for i, (title, body, color) in enumerate(features):
    row, col = divmod(i, 2)
    cx = Inches(0.4 + col * 6.4)
    cy = Inches(1.75 + row * 1.85)
    add_rect(s, cx, cy, Inches(6.1), Inches(1.7),
             fill=CARD_BG, line=color, line_w=Pt(1.5))
    add_text(s, title, cx+Inches(0.18), cy+Inches(0.15), Inches(5.7), Inches(0.4),
             size=Pt(13), bold=True, color=color)
    add_text(s, body, cx+Inches(0.18), cy+Inches(0.6), Inches(5.7), Inches(1.0),
             size=Pt(12), color=MUTED)
slide_number(s, 9, TOTAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — THANK YOU / LINKS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout); bg(s); gradient_bar(s)
add_text(s, "Thank You", Inches(1), Inches(1.4), Inches(11), Inches(1.2),
         size=Pt(64), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "TruthTrace — The AI Hallucination Auditor",
         Inches(1), Inches(2.7), Inches(11), Inches(0.7),
         size=Pt(22), color=PURPLE, align=PP_ALIGN.CENTER)
add_text(s, "Every claim. Every source. Every time.",
         Inches(1), Inches(3.4), Inches(11), Inches(0.55),
         size=Pt(16), italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Link cards
for i, (icon, label, link, color) in enumerate([
    ("🌐", "Live App", "truthtrace.vercel.app", PURPLE),
    ("📦", "GitHub", "github.com/devansh2604/truthtrace", MAGENTA),
    ("⚡", "API", "console.groq.com", GREEN),
]):
    bx = Inches(1.5 + i * 3.4)
    add_rect(s, bx, Inches(4.4), Inches(3.0), Inches(1.2),
             fill=CARD_BG, line=color, line_w=Pt(1.5))
    add_text(s, f"{icon}  {label}", bx+Inches(0.15), Inches(4.45), Inches(2.7), Inches(0.5),
             size=Pt(14), bold=True, color=color)
    add_text(s, link, bx+Inches(0.15), Inches(4.95), Inches(2.7), Inches(0.45),
             size=Pt(11), color=MUTED)

add_text(s, "Built with Next.js · Groq · Wikipedia API · DuckDuckGo · Vercel",
         Inches(1), Inches(6.5), Inches(11), Inches(0.4),
         size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)
slide_number(s, 10, TOTAL)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/devv/Desktop/TruthTrace-Presentation.pptx"
prs.save(out)
print(f"✅  Saved → {out}")
